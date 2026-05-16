"""
文档转答辩PPT生成器
支持PDF和Word文档上传，自动生成学术答辩风格PPT
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import PyPDF2
from docx import Document as DocxDocument
import re
import io
import os
from datetime import datetime

# ============================================================
# 页面配置
# ============================================================
st.set_page_config(
    page_title="答辩PPT生成器",
    page_icon="📊",
    layout="wide"
)

# ============================================================
# 文档解析模块
# ============================================================

def extract_text_from_pdf(file):
    """从PDF文件提取文本"""
    text = ""
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        st.error(f"PDF解析错误: {e}")
    return text

def extract_text_from_docx(file):
    """从Word文档提取文本"""
    text = ""
    try:
        doc = DocxDocument(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        st.error(f"Word解析错误: {e}")
    return text

def parse_document(uploaded_file):
    """解析上传的文档"""
    file_name = uploaded_file.name.lower()
    
    if file_name.endswith('.pdf'):
        return extract_text_from_pdf(uploaded_file)
    elif file_name.endswith(('.docx', '.doc')):
        return extract_text_from_docx(uploaded_file)
    else:
        st.error("不支持的文件格式，请上传PDF或Word文档")
        return ""

# ============================================================
# 内容分析模块
# ============================================================

def analyze_content(text):
    """分析文档内容，提取关键信息"""
    # 清理文本
    text = re.sub(r'\s+', ' ', text).strip()
    
    # 尝试提取标题（通常在文档开头）
    lines = text.split('\n') if '\n' in text else text.split('。')
    title = lines[0][:50] if lines else "答辩报告"
    if len(title) < 5:
        title = "学术答辩报告"
    
    # 分段处理
    sections = {
        'title': title,
        'background': "",
        'method': "",
        'result': "",
        'conclusion': "",
        'full_text': text
    }
    
    # 尝试识别常见章节
    keywords = {
        'background': ['背景', '研究背景', '问题背景', '引言', '绪论', 'background', 'introduction'],
        'method': ['方法', '研究方法', '技术路线', '方法', 'method', 'approach'],
        'result': ['结果', '实验结果', '研究结果', '分析', 'result', 'analysis'],
        'conclusion': ['结论', '总结', '展望', 'conclusion', 'summary', 'future']
    }
    
    # 简单分段提取
    text_lower = text.lower()
    for key, kws in keywords.items():
        for kw in kws:
            pattern = rf'{kw}.*?(?={("|".join([k for ks in keywords.values() for k in ks]))}|$)'
            match = re.search(pattern, text_lower, re.DOTALL)
            if match:
                sections[key] = text[match.start():match.end()][:500]
                break
    
    # 如果没有找到特定章节，按段落分配
    paragraphs = [p.strip() for p in text.split('。') if len(p.strip()) > 50]
    
    if not sections['background'] and len(paragraphs) > 0:
        sections['background'] = '。'.join(paragraphs[:max(1, len(paragraphs)//4)])
    if not sections['method'] and len(paragraphs) > 1:
        sections['method'] = '。'.join(paragraphs[max(1, len(paragraphs)//4):max(2, len(paragraphs)//2)])
    if not sections['result'] and len(paragraphs) > 2:
        sections['result'] = '。'.join(paragraphs[max(2, len(paragraphs)//2):max(3, 3*len(paragraphs)//4)])
    if not sections['conclusion'] and len(paragraphs) > 3:
        sections['conclusion'] = '。'.join(paragraphs[max(3, 3*len(paragraphs)//4):])
    
    return sections

# ============================================================
# PPT生成模块
# ============================================================

def create_defense_ppt(sections, author_name="答辩人", date_str=None):
    """生成学术答辩风格PPT"""
    
    if date_str is None:
        date_str = datetime.now().strftime("%Y年%m月")
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # 颜色定义 (学术蓝色系)
    TITLE_COLOR = RGBColor(0x1E, 0x3A, 0x5F)      # 深蓝
    ACCENT_COLOR = RGBColor(0x3A, 0x7C, 0xA5)     # 中蓝
    TEXT_COLOR = RGBColor(0x2C, 0x3E, 0x50)       # 深灰
    LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)         # 浅灰背景
    
    # 添加空白幻灯片
    def add_slide():
        blank_layout = prs.slide_layouts[6]  # 空白布局
        return prs.slides.add_slide(blank_layout)
    
    # 添加文本框
    def add_textbox(slide, left, top, width, height, text, font_size=18, 
                    bold=False, color=TEXT_COLOR, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return txBox
    
    # 添加矩形背景
    def add_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    
    # ========================================
    # 幻灯片1: 封面
    # ========================================
    slide = add_slide()
    
    # 顶部装饰条
    add_rect(slide, 0, 0, 13.333, 0.8, ACCENT_COLOR)
    
    # 标题
    title = sections['title'][:40] + "..." if len(sections['title']) > 40 else sections['title']
    add_textbox(slide, 1, 2.5, 11.333, 1.5, title, 
                font_size=44, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    
    # 副标题
    add_textbox(slide, 1, 4.2, 11.333, 0.6, "学术答辩报告",
                font_size=24, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
    
    # 答辩人信息
    add_textbox(slide, 1, 5.5, 11.333, 0.5, f"答辩人：{author_name}",
                font_size=18, color=TEXT_COLOR, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 6.0, 11.333, 0.5, date_str,
                font_size=16, color=TEXT_COLOR, align=PP_ALIGN.CENTER)
    
    # 底部装饰条
    add_rect(slide, 0, 7.1, 13.333, 0.4, ACCENT_COLOR)
    
    # ========================================
    # 幻灯片2: 目录
    # ========================================
    slide = add_slide()
    
    # 标题
    add_textbox(slide, 0.5, 0.5, 12.333, 0.8, "目 录",
                font_size=36, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    
    # 目录项
    toc_items = [
        "01  研究背景",
        "02  研究方法",
        "03  研究结果",
        "04  结论与展望",
        "05  致谢"
    ]
    
    for i, item in enumerate(toc_items):
        y_pos = 1.8 + i * 1.0
        # 编号背景
        add_rect(slide, 1.5, y_pos, 0.6, 0.6, ACCENT_COLOR)
        # 目录文字
        add_textbox(slide, 2.3, y_pos + 0.1, 8, 0.5, item[3:],
                    font_size=24, color=TEXT_COLOR)
    
    # ========================================
    # 幻灯片3: 研究背景
    # ========================================
    slide = add_slide()
    
    # 章节标题背景
    add_rect(slide, 0, 0, 13.333, 1.2, TITLE_COLOR)
    add_textbox(slide, 0.5, 0.35, 12.333, 0.6, "01  研究背景",
                font_size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    
    # 内容
    content = sections['background'][:600] if sections['background'] else "（请在此补充研究背景内容）"
    # 分段显示
    paragraphs = content.split('。')[:5]
    for i, para in enumerate(paragraphs):
        if para.strip():
            add_textbox(slide, 0.8, 1.6 + i * 1.0, 11.5, 0.9, 
                       f"• {para.strip()}。",
                       font_size=18, color=TEXT_COLOR)
    
    # ========================================
    # 幻灯片4: 研究方法
    # ========================================
    slide = add_slide()
    
    add_rect(slide, 0, 0, 13.333, 1.2, TITLE_COLOR)
    add_textbox(slide, 0.5, 0.35, 12.333, 0.6, "02  研究方法",
                font_size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    
    content = sections['method'][:600] if sections['method'] else "（请在此补充研究方法内容）"
    paragraphs = content.split('。')[:5]
    for i, para in enumerate(paragraphs):
        if para.strip():
            add_textbox(slide, 0.8, 1.6 + i * 1.0, 11.5, 0.9,
                       f"• {para.strip()}。",
                       font_size=18, color=TEXT_COLOR)
    
    # ========================================
    # 幻灯片5: 研究结果
    # ========================================
    slide = add_slide()
    
    add_rect(slide, 0, 0, 13.333, 1.2, TITLE_COLOR)
    add_textbox(slide, 0.5, 0.35, 12.333, 0.6, "03  研究结果",
                font_size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    
    content = sections['result'][:600] if sections['result'] else "（请在此补充研究结果内容）"
    paragraphs = content.split('。')[:5]
    for i, para in enumerate(paragraphs):
        if para.strip():
            add_textbox(slide, 0.8, 1.6 + i * 1.0, 11.5, 0.9,
                       f"• {para.strip()}。",
                       font_size=18, color=TEXT_COLOR)
    
    # ========================================
    # 幻灯片6: 结论与展望
    # ========================================
    slide = add_slide()
    
    add_rect(slide, 0, 0, 13.333, 1.2, TITLE_COLOR)
    add_textbox(slide, 0.5, 0.35, 12.333, 0.6, "04  结论与展望",
                font_size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    
    content = sections['conclusion'][:600] if sections['conclusion'] else "（请在此补充结论与展望内容）"
    paragraphs = content.split('。')[:5]
    for i, para in enumerate(paragraphs):
        if para.strip():
            add_textbox(slide, 0.8, 1.6 + i * 1.0, 11.5, 0.9,
                       f"• {para.strip()}。",
                       font_size=18, color=TEXT_COLOR)
    
    # ========================================
    # 幻灯片7: 致谢
    # ========================================
    slide = add_slide()
    
    # 全屏背景
    add_rect(slide, 0, 0, 13.333, 7.5, TITLE_COLOR)
    
    add_textbox(slide, 1, 2.8, 11.333, 1.0, "感谢各位老师指导！",
                font_size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), 
                align=PP_ALIGN.CENTER)
    
    add_textbox(slide, 1, 4.2, 11.333, 0.6, "Thank You",
                font_size=28, color=RGBColor(0xFF, 0xFF, 0xFF), 
                align=PP_ALIGN.CENTER)
    
    add_textbox(slide, 1, 5.5, 11.333, 0.5, f"答辩人：{author_name}",
                font_size=18, color=RGBColor(0xFF, 0xFF, 0xFF), 
                align=PP_ALIGN.CENTER)
    
    return prs

# ============================================================
# 主界面
# ============================================================

def main():
    # 标题
    st.title("📊 答辩PPT生成器")
    st.markdown("---")
    
    # 侧边栏设置
    with st.sidebar:
        st.header("⚙️ 设置")
        author_name = st.text_input("答辩人姓名", value="答辩人")
        date_str = st.text_input("答辩日期", value=datetime.now().strftime("%Y年%m月"))
        
        st.markdown("---")
        st.markdown("### 📝 使用说明")
        st.markdown("""
        1. 上传PDF或Word文档
        2. 系统自动分析文档内容
        3. 生成学术答辩风格PPT
        4. 下载并编辑完善
        """)
    
    # 文件上传
    col1, col2 = st.columns([2, 1])
    
    with col1:
        uploaded_file = st.file_uploader(
            "上传文档",
            type=['pdf', 'docx', 'doc'],
            help="支持PDF和Word文档格式"
        )
    
    with col2:
        st.markdown("### 支持格式")
        st.markdown("- 📄 PDF文档")
        st.markdown("- 📝 Word文档")
    
    if uploaded_file is not None:
        st.success(f"已上传: {uploaded_file.name}")
        
        # 解析文档
        with st.spinner("正在解析文档..."):
            text = parse_document(uploaded_file)
        
        if text:
            # 显示提取的文本预览
            with st.expander("📖 查看提取内容", expanded=False):
                st.text_area("文档内容", text[:2000] + "..." if len(text) > 2000 else text, 
                            height=200)
            
            # 分析内容
            with st.spinner("正在分析内容..."):
                sections = analyze_content(text)
            
            # 显示分析结果
            col_a, col_b = st.columns(2)
            with col_a:
                st.markdown("**📌 识别到的标题:**")
                st.info(sections['title'][:50])
            
            with col_b:
                st.markdown("**📊 内容统计:**")
                st.info(f"总字数: {len(text)} 字")
            
            # 生成PPT按钮
            if st.button("🎬 生成答辩PPT", type="primary"):
                with st.spinner("正在生成PPT..."):
                    prs = create_defense_ppt(sections, author_name, date_str)
                    
                    # 保存到内存
                    ppt_buffer = io.BytesIO()
                    prs.save(ppt_buffer)
                    ppt_buffer.seek(0)
                    
                    # 生成文件名
                    output_name = uploaded_file.name.rsplit('.', 1)[0] + "_答辩PPT.pptx"
                    
                    # 下载按钮
                    st.success("✅ PPT生成成功！")
                    st.download_button(
                        label="📥 下载PPT文件",
                        data=ppt_buffer,
                        file_name=output_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
                    # PPT结构说明
                    st.markdown("### 📋 PPT结构")
                    st.markdown("""
                    | 序号 | 页面 | 内容 |
                    |:----:|:----:|:-----|
                    | 1 | 封面 | 标题、答辩人、日期 |
                    | 2 | 目录 | 五大章节导航 |
                    | 3 | 研究背景 | 文档背景部分 |
                    | 4 | 研究方法 | 文档方法部分 |
                    | 5 | 研究结果 | 文档结果部分 |
                    | 6 | 结论与展望 | 文档结论部分 |
                    | 7 | 致谢 | 结束页 |
                    """)

if __name__ == "__main__":
    main()
